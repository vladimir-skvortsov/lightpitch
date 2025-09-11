"""
Presentation Generator using LLM to create improved presentations with visual suggestions
"""

import json
import logging
import os
import sys
import time
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Add the project root to the path to import other modules
sys.path.append(os.path.join(os.path.dirname(__file__), '..', '..'))

from models.text_editor.openrouter_client import OpenRouterService
from models.presentation_summary.presentation_summarizer import PresentationSummarizer
from .types import ImprovedSlide, ImprovedPresentation, VisualElement, VisualElementType, ChartType

logger = logging.getLogger(__name__)


class PresentationGenerator:
    """
    Generates improved presentations with visual element suggestions using LLM
    """
    
    def __init__(self, model: str = 'anthropic/claude-3.5-haiku'):
        """Initialize with OpenRouter service and presentation summarizer"""
        self.openai_service = OpenRouterService(model=model)
        self.summarizer = PresentationSummarizer(model=model)
        
        # Improvement prompt with visual suggestions
        self.improvement_prompt = """
Ты эксперт по созданию презентаций. На основе анализа существующей презентации, создай улучшенную версию с предложениями визуальных элементов.

ПРАВИЛА УЛУЧШЕНИЯ:
1. **Убирай канцеляризмы и лишние слова**: "в принципе", "как бы", "собственно", "значит", "ну", "в общем"
2. **Исправляй грамматику и стилистику**: правильные падежи, согласование, логичные переходы
3. **Делай тексты более четкими**: одна мысль на слайд, конкретные формулировки
4. **Улучшай структуру**: логичная последовательность, понятные заголовки
5. **Добавляй call-to-action**: где это уместно
6. **Предлагай визуальные элементы**: диаграммы, графики, изображения, таблицы, иконки

ТЕКУЩАЯ ПРЕЗЕНТАЦИЯ:
{original_content}

АНАЛИЗ И РЕКОМЕНДАЦИИ:
{warnings_errors}

Верни результат в JSON формате:
{{
  "improved_slides": [
    {{
      "slide_number": номер_слайда,
      "title": "улучшенный_заголовок",
      "content": ["улучшенный_текст_1", "улучшенный_текст_2"],
      "bullet_points": ["пункт_1", "пункт_2", "пункт_3"],
      "speaker_notes": "заметки_для_докладчика",
      "improvements_applied": ["что_было_улучшено"],
      "suggested_visuals": [
        {{
          "element_type": "chart|graph|diagram|image|table|icon|infografic|timeline|flowchart",
          "title": "название_элемента",
          "description": "описание_элемента",
          "purpose": "зачем_нужен_этот_элемент",
          "data_suggestion": ["данные_для_диаграммы"],
          "chart_type": "bar|line|pie|column|area|scatter" (только для charts),
          "position_suggestion": "center|left|right|top|bottom",
          "size_suggestion": "small|medium|large"
        }}
      ]
    }}
  ],
  "improvements_summary": [
    "список_основных_улучшений"
  ],
  "overall_feedback": "общая_оценка_улучшений"
}}

ПРАВИЛА ДЛЯ ВИЗУАЛЬНЫХ ЭЛЕМЕНТОВ:
- Предлагай диаграммы для числовых данных и сравнений
- Используй графики для показа трендов и изменений во времени
- Добавляй диаграммы для процессов и схем
- Предлагай таблицы для структурированных данных
- Рекомендуй иконки для ключевых понятий
- Используй инфографику для сложных концепций
- Добавляй timeline для хронологических данных

Создай профессиональную, четкую и визуально привлекательную презентацию.
"""
    
    def extract_analysis_issues(self, analysis: Dict[str, Any]) -> str:
        """Extract warnings and errors from analysis for improvement prompt"""
        issues = []
        
        # Add warnings
        for warning in analysis.get('warnings', []):
            issues.append(f"⚠️ {warning.get('title', '')}: {warning.get('description', '')}")
        
        # Add errors
        for error in analysis.get('errors', []):
            issues.append(f"❌ {error.get('title', '')}: {error.get('description', '')}")
        
        # Add recommendations
        for rec in analysis.get('recommendations', []):
            issues.append(f"💡 Рекомендация: {rec}")
        
        return "\n".join(issues) if issues else "Особых проблем не обнаружено"
    
    def parse_visual_element(self, visual_data: Dict[str, Any]) -> VisualElement:
        """Parse visual element from LLM response"""
        try:
            element_type = VisualElementType(visual_data.get('element_type', 'chart'))
            chart_type = None
            
            if element_type in [VisualElementType.CHART, VisualElementType.GRAPH]:
                chart_type_str = visual_data.get('chart_type')
                if chart_type_str:
                    chart_type = ChartType(chart_type_str)
            
            return VisualElement(
                element_type=element_type,
                title=visual_data.get('title', ''),
                description=visual_data.get('description', ''),
                purpose=visual_data.get('purpose', ''),
                data_suggestion=visual_data.get('data_suggestion', []),
                chart_type=chart_type,
                position_suggestion=visual_data.get('position_suggestion', 'center'),
                size_suggestion=visual_data.get('size_suggestion', 'medium')
            )
        except (ValueError, KeyError) as e:
            logger.warning(f"Error parsing visual element: {str(e)}")
            # Return default visual element
            return VisualElement(
                element_type=VisualElementType.CHART,
                title=visual_data.get('title', 'Диаграмма'),
                description=visual_data.get('description', 'Визуальный элемент'),
                purpose=visual_data.get('purpose', 'Иллюстрация данных')
            )
    
    async def improve_slide_content(self, original_content: str, analysis_issues: str) -> Dict[str, Any]:
        """Improve slide content using LLM with visual suggestions"""
        try:
            prompt = self.improvement_prompt.format(
                original_content=original_content,
                warnings_errors=analysis_issues
            )
            
            logger.info("Improving presentation content with visual suggestions")
            improved_json = await self.openai_service.analyze_text(
                prompt=prompt,
                text="",
                expect_json=True
            )
            
            improved_data = json.loads(improved_json)
            return improved_data
            
        except Exception as e:
            logger.error(f"Error improving slide content: {str(e)}")
            return {
                'improved_slides': [],
                'improvements_summary': [f'Ошибка улучшения: {str(e)}'],
                'overall_feedback': 'Не удалось улучшить презентацию'
            }
    
    def create_improved_pptx(self, improved_data: Dict[str, Any], original_file_path: str, output_path: str) -> str:
        """Create improved .pptx file with visual placeholders"""
        try:
            # Create new presentation
            prs = Presentation()
            
            # Set slide size to standard (16:9)
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            
            improved_slides = improved_data.get('improved_slides', [])
            
            for slide_data in improved_slides:
                # Add new slide
                slide_layout = prs.slide_layouts[1]  # Title and Content layout
                slide = prs.slides.add_slide(slide_layout)
                
                # Add title
                title_shape = slide.shapes.title
                title_shape.text = slide_data.get('title', 'Заголовок слайда')
                
                # Add content to content placeholder
                content_shape = slide.placeholders[1]
                tf = content_shape.text_frame
                tf.clear()
                
                # Add main content
                content_items = slide_data.get('content', [])
                for i, content in enumerate(content_items):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    p.text = content
                    p.font.size = Pt(18)
                    p.font.color.rgb = RGBColor(51, 51, 51)
                
                # Add bullet points
                bullet_points = slide_data.get('bullet_points', [])
                if bullet_points:
                    if content_items:  # Add spacing if there's already content
                        p = tf.add_paragraph()
                        p.text = ""
                    
                    for bullet in bullet_points:
                        p = tf.add_paragraph()
                        p.text = bullet
                        p.font.size = Pt(16)
                        p.font.color.rgb = RGBColor(51, 51, 51)
                        p.level = 0
                
                # Add visual suggestions as text placeholders
                suggested_visuals = slide_data.get('suggested_visuals', [])
                if suggested_visuals:
                    p = tf.add_paragraph()
                    p.text = ""
                    
                    p = tf.add_paragraph()
                    p.text = "РЕКОМЕНДУЕМЫЕ ВИЗУАЛЬНЫЕ ЭЛЕМЕНТЫ:"
                    p.font.size = Pt(14)
                    p.font.color.rgb = RGBColor(0, 100, 0)
                    p.font.bold = True
                    
                    for visual in suggested_visuals:
                        p = tf.add_paragraph()
                        p.text = f"• {visual.get('title', '')}: {visual.get('description', '')}"
                        p.font.size = Pt(12)
                        p.font.color.rgb = RGBColor(0, 100, 0)
                        p.level = 1
                
                # Add speaker notes if available
                speaker_notes = slide_data.get('speaker_notes')
                if speaker_notes:
                    notes_slide = slide.notes_slide
                    notes_text_frame = notes_slide.notes_text_frame
                    notes_text_frame.text = speaker_notes
            
            # Save presentation
            prs.save(output_path)
            logger.info(f"Improved presentation with visual suggestions saved to: {output_path}")
            return output_path
            
        except Exception as e:
            logger.error(f"Error creating improved .pptx: {str(e)}")
            raise
    
    async def generate_improved_presentation(self, original_file_path: str, analysis: Dict[str, Any], output_dir: str = None) -> ImprovedPresentation:
        """
        Generate improved presentation with visual suggestions
        
        Args:
            original_file_path: Path to original .pptx file
            analysis: Analysis results from presentation_summary
            output_dir: Directory to save improved presentation
            
        Returns:
            ImprovedPresentation object with visual suggestions
        """
        try:
            logger.info(f"Generating improved presentation with visual suggestions for: {original_file_path}")
            
            # Parse original presentation
            original_content = self.summarizer.parse_presentation(original_file_path)
            
            # Format original content for LLM
            formatted_content = self.summarizer.format_presentation_for_analysis(original_content)
            
            # Extract analysis issues
            analysis_issues = self.extract_analysis_issues(analysis)
            
            # Improve content with LLM
            improved_data = await self.improve_slide_content(formatted_content, analysis_issues)
            
            # Create improved slides objects with visual elements
            improved_slides = []
            for slide_data in improved_data.get('improved_slides', []):
                # Parse visual elements
                suggested_visuals = []
                for visual_data in slide_data.get('suggested_visuals', []):
                    visual_element = self.parse_visual_element(visual_data)
                    suggested_visuals.append(visual_element)
                
                improved_slide = ImprovedSlide(
                    slide_number=slide_data.get('slide_number', 0),
                    title=slide_data.get('title', ''),
                    content=slide_data.get('content', []),
                    bullet_points=slide_data.get('bullet_points', []),
                    speaker_notes=slide_data.get('speaker_notes'),
                    improvements_applied=slide_data.get('improvements_applied', []),
                    suggested_visuals=suggested_visuals
                )
                improved_slides.append(improved_slide)
            
            # Generate output filename
            original_filename = os.path.basename(original_file_path)
            name, ext = os.path.splitext(original_filename)
            improved_filename = f"{name}_improved{ext}"
            
            # Set output directory
            if output_dir is None:
                output_dir = os.path.dirname(original_file_path)
            
            output_path = os.path.join(output_dir, improved_filename)
            
            # Create improved .pptx file
            self.create_improved_pptx(improved_data, original_file_path, output_path)
            
            # Create improved presentation object
            improved_presentation = ImprovedPresentation(
                original_filename=original_filename,
                improved_filename=improved_filename,
                total_slides=len(improved_slides),
                slides=improved_slides,
                improvements_summary=improved_data.get('improvements_summary', []),
                generation_timestamp=str(int(time.time()))
            )
            
            logger.info(f"Successfully generated improved presentation with visual suggestions: {improved_filename}")
            return improved_presentation
            
        except Exception as e:
            logger.error(f"Error generating improved presentation: {str(e)}")
            raise


async def generate_improved_presentation(original_file_path: str, analysis: Dict[str, Any], output_dir: str = None, model: str = 'anthropic/claude-3.5-haiku') -> ImprovedPresentation:
    """
    Convenience function to generate improved presentation with visual suggestions
    
    Args:
        original_file_path: Path to original .pptx file
        analysis: Analysis results
        output_dir: Output directory
        model: LLM model to use
        
    Returns:
        ImprovedPresentation object with visual suggestions
    """
    try:
        generator = PresentationGenerator(model=model)
        return await generator.generate_improved_presentation(original_file_path, analysis, output_dir)
    except Exception as e:
        logger.error(f"Presentation generation failed: {str(e)}")
        raise
