"""
Presentation Summarizer using LLM for comprehensive presentation analysis
"""

import json
import logging
import os
import sys
from typing import Dict, List, Any, Optional
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Add the project root to the path to import openrouter_client
sys.path.append(os.path.join(os.path.dirname(__file__), '..', '..'))

from models.text_editor.openrouter_client import OpenRouterService
from .types import SlideContent, PresentationContent, AnalysisResult

logger = logging.getLogger(__name__)


class PresentationSummarizer:
    """
    Analyzes presentation files (.pptx) using LLM to provide comprehensive feedback
    """
    
    def __init__(self, model: str = 'anthropic/claude-3.5-haiku'):
        """Initialize with OpenRouter service"""
        self.openai_service = OpenRouterService(model=model)
        
        # Analysis prompts
        self.analysis_prompt = """
Ты эксперт по презентациям и дизайну. Проанализируй презентацию по следующим критериям:

1. **Структура и логика**:
   - Логическая последовательность слайдов
   - Ясность повествования
   - Соответствие цели презентации

2. **Контент и тексты**:
   - Качество текстов (орфография, пунктуация, стилистика)
   - Соответствие принципу "одна мысль на слайд"
   - Удаление лишних слов и канцеляризмов
   - Четкость формулировок

3. **Визуальное оформление**:
   - Консистентность дизайна
   - Читаемость шрифтов
   - Цветовая схема
   - Использование пространства

4. **Профессионализм**:
   - Корпоративный стиль
   - Качество исполнения
   - Внимание к деталям

Верни результат в JSON формате:
{
  "overall_score": число от 0 до 100,
  "category_scores": {
    "structure": число от 0 до 10,
    "content": число от 0 до 10,
    "design": число от 0 до 10,
    "professionalism": число от 0 до 10
  },
  "good_practices": [
    {
      "title": "название практики",
      "description": "описание",
      "category": "категория"
    }
  ],
  "warnings": [
    {
      "title": "название предупреждения",
      "description": "описание проблемы",
      "category": "категория",
      "slides": [номера_слайдов]
    }
  ],
  "errors": [
    {
      "title": "название ошибки",
      "description": "описание критической проблемы",
      "category": "категория",
      "slides": [номера_слайдов],
      "severity": "high/medium/low"
    }
  ],
  "recommendations": ["список рекомендаций"],
  "strengths": ["список сильных сторон"],
  "areas_for_improvement": ["список областей для улучшения"],
  "feedback": "общая оценка презентации"
}

Будь критичным, но конструктивным. Давай конкретные рекомендации для улучшения.
"""
    
    def extract_text_from_shape(self, shape) -> List[str]:
        """Extract text content from a shape"""
        texts = []
        
        if hasattr(shape, 'text') and shape.text.strip():
            texts.append(shape.text.strip())
        
        # Handle text in table cells
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
        
        return texts
    
    def extract_slide_content(self, slide, slide_number: int) -> SlideContent:
        """Extract content from a single slide"""
        title = None
        text_content = []
        bullet_points = []
        speaker_notes = None
        
        # Extract speaker notes
        if hasattr(slide, 'notes_slide') and slide.notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text
            if notes_text.strip():
                speaker_notes = notes_text.strip()
        
        # Extract content from shapes
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                
                # Try to identify title (usually first text or largest font)
                if not title and len(text) < 100:
                    title = text
                else:
                    # Check if it's bullet points (contains bullet characters or is in a list)
                    if any(char in text for char in ['•', '◦', '-', '*']) or '\n' in text:
                        bullet_points.extend([line.strip() for line in text.split('\n') if line.strip()])
                    else:
                        text_content.append(text)
        
        return SlideContent(
            slide_number=slide_number,
            title=title,
            text_content=text_content,
            bullet_points=bullet_points,
            speaker_notes=speaker_notes
        )
    
    def parse_presentation(self, file_path: str) -> PresentationContent:
        """Parse .pptx file and extract content"""
        try:
            prs = Presentation(file_path)
            
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                slide_content = self.extract_slide_content(slide, i)
                slides.append(slide_content)
            
            return PresentationContent(
                filename=os.path.basename(file_path),
                total_slides=len(slides),
                slides=slides
            )
            
        except Exception as e:
            logger.error(f"Error parsing presentation {file_path}: {str(e)}")
            raise
    
    def format_presentation_for_analysis(self, content: PresentationContent) -> str:
        """Format presentation content for LLM analysis"""
        formatted_text = f"Презентация: {content.filename}\n"
        formatted_text += f"Общее количество слайдов: {content.total_slides}\n\n"
        
        for slide in content.slides:
            formatted_text += f"--- Слайд {slide.slide_number} ---\n"
            
            if slide.title:
                formatted_text += f"Заголовок: {slide.title}\n"
            
            if slide.text_content:
                formatted_text += "Основной текст:\n"
                for text in slide.text_content:
                    formatted_text += f"- {text}\n"
            
            if slide.bullet_points:
                formatted_text += "Список/маркеры:\n"
                for point in slide.bullet_points:
                    formatted_text += f"• {point}\n"
            
            if slide.speaker_notes:
                formatted_text += f"Заметки докладчика: {slide.speaker_notes}\n"
            
            formatted_text += "\n"
        
        return formatted_text
    
    async def analyze_presentation_content(self, content: PresentationContent) -> Dict[str, Any]:
        """Analyze presentation content using LLM"""
        try:
            formatted_text = self.format_presentation_for_analysis(content)
            
            logger.info(f"Analyzing presentation: {content.filename}")
            analysis_json = await self.openai_service.analyze_text(
                prompt=self.analysis_prompt,
                text=formatted_text,
                expect_json=True
            )
            
            analysis = json.loads(analysis_json)
            
            # Ensure all required fields are present
            if 'overall_score' not in analysis:
                analysis['overall_score'] = 0
            if 'category_scores' not in analysis:
                analysis['category_scores'] = {}
            if 'good_practices' not in analysis:
                analysis['good_practices'] = []
            if 'warnings' not in analysis:
                analysis['warnings'] = []
            if 'errors' not in analysis:
                analysis['errors'] = []
            if 'recommendations' not in analysis:
                analysis['recommendations'] = []
            if 'strengths' not in analysis:
                analysis['strengths'] = []
            if 'areas_for_improvement' not in analysis:
                analysis['areas_for_improvement'] = []
            
            return analysis
            
        except Exception as e:
            logger.error(f"Error analyzing presentation content: {str(e)}")
            return {
                'error': f'Analysis failed: {str(e)}',
                'overall_score': 0,
                'category_scores': {},
                'good_practices': [],
                'warnings': [],
                'errors': [],
                'recommendations': [f'Ошибка анализа: {str(e)}'],
                'strengths': [],
                'areas_for_improvement': [],
            }
    
    async def analyze_presentation_file(self, file_path: str) -> Dict[str, Any]:
        """
        Analyze a presentation file and return comprehensive analysis
        
        Args:
            file_path: Path to the .pptx file
            
        Returns:
            Dictionary with analysis results
        """
        try:
            # Parse presentation
            content = self.parse_presentation(file_path)
            
            # Analyze with LLM
            analysis = await self.analyze_presentation_content(content)
            
            # Add metadata
            analysis['filename'] = content.filename
            analysis['total_slides'] = content.total_slides
            analysis['analysis_method'] = 'LLM анализ'
            analysis['analysis_timestamp'] = str(int(__import__('time').time()))
            
            return analysis
            
        except Exception as e:
            logger.error(f"Presentation analysis failed: {str(e)}")
            return {
                'error': f'Analysis failed: {str(e)}',
                'overall_score': 0,
                'filename': os.path.basename(file_path),
                'total_slides': 0,
                'category_scores': {},
                'good_practices': [],
                'warnings': [],
                'errors': [],
                'recommendations': [f'Ошибка анализа: {str(e)}'],
                'strengths': [],
                'areas_for_improvement': [],
                'analysis_method': 'LLM анализ',
            }


async def analyze_presentation(file_path: str, model: str = 'anthropic/claude-3.5-haiku') -> Dict[str, Any]:
    """
    Convenience function to analyze a presentation file
    
    Args:
        file_path: Path to the .pptx file
        model: LLM model to use
        
    Returns:
        Analysis results dictionary
    """
    try:
        summarizer = PresentationSummarizer(model=model)
        return await summarizer.analyze_presentation_file(file_path)
    except Exception as e:
        logger.error(f"Presentation analysis failed: {str(e)}")
        return {
            'error': f'Analysis failed: {str(e)}',
            'overall_score': 0,
            'filename': os.path.basename(file_path),
            'total_slides': 0,
            'category_scores': {},
            'good_practices': [],
            'warnings': [],
            'errors': [],
            'recommendations': [f'Ошибка анализа: {str(e)}'],
            'strengths': [],
            'areas_for_improvement': [],
        }
